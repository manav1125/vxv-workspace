from __future__ import annotations

import os
from typing import Optional
from uuid import uuid4
from io import BytesIO

from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse

from .models import ModuleKey
from .models import (
    AgentProfile,
    AgentUpdateRequest,
    ActionResponse,
    AppLaunchRequest,
    AuthSession,
    ApprovalRequest,
    Artifact,
    ArtifactUpdateRequest,
    BootstrapResponse,
    ChatRequest,
    ChatResponse,
    Contact,
    ContactCreateRequest,
    FundraiseInvestor,
    FundraiseInvestorCreateRequest,
    FundraiseInvestorUpdateRequest,
    Goal,
    GoalCreateRequest,
    GoalUpdateRequest,
    InvestorRoomActionResponse,
    KnowledgeSource,
    KnowledgeSourceCreateRequest,
    LoginRequest,
    PublishInvestorRoomRequest,
    UploadResponse,
    UploadRecord,
    Workspace,
    WorkspaceSetupRequest,
    WorkspaceUser,
    WorkspaceUserCreateRequest,
    WorkspaceUserUpdateRequest,
    WorkflowLaunchRequest,
    now_iso,
)
from .auth import OIDCVerifier
from .orchestrator import FounderOrchestrator
from .persistence import PersistenceBackend
from .storage import UploadStorage
from .store import DemoStore

app = FastAPI(
    title="VXV Workspace API",
    version="0.1.0",
    summary="Unified founder OS scaffold aligned to the AgentScope ecosystem.",
)

cors_origins_env = os.getenv("CORS_ORIGINS", "*").strip()
cors_origins = ["*"] if cors_origins_env in {"", "*"} else [
    origin.strip() for origin in cors_origins_env.split(",") if origin.strip()
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=cors_origins,
    allow_credentials=cors_origins != ["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

store = DemoStore()
orchestrator = FounderOrchestrator(store)
auth_backend = PersistenceBackend()
oidc_verifier = OIDCVerifier()
upload_storage = UploadStorage()


@app.middleware("http")
async def auth_middleware(request: Request, call_next):
    public_paths = {"/health", "/api/auth/login"}
    if request.method == "OPTIONS" or request.url.path in public_paths or not request.url.path.startswith("/api/"):
        return await call_next(request)

    session = resolve_session(request)
    if session is None:
        return JSONResponse(status_code=401, content={"detail": "Authentication required"})

    request.state.session = session
    return await call_next(request)


@app.get("/health")
def health() -> dict[str, str]:
    return {"status": "ok"}


@app.get("/api/bootstrap", response_model=BootstrapResponse)
def get_bootstrap() -> BootstrapResponse:
    return store.bootstrap()


@app.post("/api/auth/login", response_model=AuthSession)
def post_login(request: LoginRequest) -> AuthSession:
    session = auth_backend.create_session(request.email, request.password)
    if session is None:
        raise HTTPException(status_code=401, detail="Invalid email or password")
    return AuthSession.model_validate(session)


@app.get("/api/auth/session", response_model=AuthSession)
def get_session(request: Request) -> AuthSession:
    session = getattr(request.state, "session", None)
    if session is None:
        raise HTTPException(status_code=401, detail="Authentication required")
    return AuthSession.model_validate(session)


@app.get("/api/users", response_model=list[WorkspaceUser])
def get_users(request: Request) -> list[WorkspaceUser]:
    session = request.state.session
    return [WorkspaceUser.model_validate(item) for item in auth_backend.list_users(session["workspace_id"])]


@app.post("/api/users", response_model=WorkspaceUser)
def post_user(request: Request, payload: WorkspaceUserCreateRequest) -> WorkspaceUser:
    session = request.state.session
    if session.get("role") != "owner":
        raise HTTPException(status_code=403, detail="Only owners can create workspace users")
    user = auth_backend.create_user(
        email=payload.email,
        password=payload.password,
        workspace_id=session["workspace_id"],
        display_name=payload.display_name,
        role=payload.role,
    )
    return WorkspaceUser.model_validate(user)


@app.patch("/api/users/{user_email:path}", response_model=WorkspaceUser)
def patch_user(request: Request, user_email: str, payload: WorkspaceUserUpdateRequest) -> WorkspaceUser:
    session = request.state.session
    if session.get("role") != "owner":
        raise HTTPException(status_code=403, detail="Only owners can manage workspace users")
    try:
        user = auth_backend.update_user(
            email=user_email.lower(),
            workspace_id=session["workspace_id"],
            display_name=payload.display_name,
            role=payload.role,
            status=payload.status,
            password=payload.password,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc
    return WorkspaceUser.model_validate(user)


@app.get("/api/workspaces/current")
def get_workspace() -> Workspace:
    return store.workspace


@app.patch("/api/workspaces/current", response_model=Workspace)
def patch_workspace(request: WorkspaceSetupRequest) -> Workspace:
    return store.update_workspace(
        company_name=request.company_name,
        founder_name=request.founder_name,
        stage=request.stage,
        mission=request.mission,
        primary_kpi=request.primary_kpi,
        summary=request.summary,
    )


@app.get("/api/goals")
def get_goals() -> list[Goal]:
    return store.goals


@app.post("/api/goals", response_model=Goal)
def post_goal(request: GoalCreateRequest) -> Goal:
    return store.add_goal(
        title=request.title,
        owner=request.owner,
        kpi=request.kpi,
        due_date=request.due_date,
        linked_agents=request.linked_agents,
        status=request.status,
    )


@app.patch("/api/goals/{goal_id}", response_model=Goal)
def patch_goal(goal_id: str, request: GoalUpdateRequest) -> Goal:
    try:
        return store.update_goal_status(goal_id, request.status)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.get("/api/agents")
def get_agents() -> list[AgentProfile]:
    return store.agents


@app.patch("/api/agents/{agent_id}", response_model=AgentProfile)
def patch_agent(agent_id: str, request: AgentUpdateRequest) -> AgentProfile:
    try:
        return store.update_agent(
            agent_id=agent_id,
            budget=request.budget,
            permissions=request.permissions,
            escalation_rule=request.escalation_rule,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.get("/api/skills")
def get_skills():
    return store.skills


@app.get("/api/apps")
def get_apps():
    return store.apps


@app.get("/api/knowledge-sources")
def get_knowledge_sources() -> list[KnowledgeSource]:
    return store.knowledge_sources


@app.post("/api/knowledge-sources", response_model=KnowledgeSource)
def post_knowledge_source(request: KnowledgeSourceCreateRequest) -> KnowledgeSource:
    return store.add_knowledge_source(
        title=request.title,
        source_type=request.source_type,
        status=request.status,
        freshness=request.freshness,
    )


@app.get("/api/task-runs")
def get_task_runs():
    return store.task_runs


@app.get("/api/artifacts")
def get_artifacts():
    return store.artifacts


@app.patch("/api/artifacts/{artifact_id}", response_model=Artifact)
def patch_artifact(artifact_id: str, request: ArtifactUpdateRequest) -> Artifact:
    try:
        return store.update_artifact_content(artifact_id, request.content)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.get("/api/contacts")
def get_contacts() -> list[Contact]:
    return store.contacts


@app.post("/api/contacts", response_model=Contact)
def post_contact(request: ContactCreateRequest) -> Contact:
    return store.add_contact(
        name=request.name,
        category=request.category,
        company=request.company,
        relationship_stage=request.relationship_stage,
        last_touch=request.last_touch or request_last_touch_today(),
    )


@app.get("/api/workflows")
def get_workflows():
    return store.workflows


@app.post("/api/workflows/{workflow_id}/launch", response_model=ActionResponse)
def post_workflow_launch(workflow_id: str, request: WorkflowLaunchRequest) -> ActionResponse:
    try:
        return orchestrator.launch_workflow(workflow_id, request)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.get("/api/fundraise-pipeline/current")
def get_fundraise_pipeline():
    return store.fundraise_pipeline


@app.post("/api/fundraise-pipeline/current/investors", response_model=FundraiseInvestor)
def post_fundraise_investor(request: FundraiseInvestorCreateRequest) -> FundraiseInvestor:
    return store.add_fundraise_investor(
        name=request.name,
        thesis=request.thesis,
        stage_fit=request.stage_fit,
        relationship_status=request.relationship_status,
        next_step=request.next_step,
    )


@app.patch("/api/fundraise-pipeline/current/investors/{investor_id}", response_model=FundraiseInvestor)
def patch_fundraise_investor(investor_id: str, request: FundraiseInvestorUpdateRequest) -> FundraiseInvestor:
    try:
        return store.update_fundraise_investor(
            investor_id=investor_id,
            relationship_status=request.relationship_status,
            next_step=request.next_step,
        )
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.get("/api/investor-room/current")
def get_investor_room():
    return store.investor_room


@app.post("/api/task-runs/{task_id}/approval", response_model=ActionResponse)
def post_task_approval(task_id: str, request: ApprovalRequest) -> ActionResponse:
    try:
        return orchestrator.decide_approval(task_id, request.decision)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.post("/api/apps/{app_id}/launch", response_model=ActionResponse)
def post_app_launch(app_id: str, request: AppLaunchRequest) -> ActionResponse:
    try:
        return orchestrator.launch_app(app_id, request)
    except KeyError as exc:
        raise HTTPException(status_code=404, detail=str(exc)) from exc


@app.post("/api/chat", response_model=ChatResponse)
def post_chat(request: ChatRequest) -> ChatResponse:
    return orchestrator.respond(request)


@app.post("/api/investor-room/publish", response_model=InvestorRoomActionResponse)
def post_publish_investor_room(request: PublishInvestorRoomRequest) -> InvestorRoomActionResponse:
    return orchestrator.publish_investor_room(request)


@app.post("/api/uploads", response_model=UploadResponse)
async def post_upload(
    request: Request,
    file: UploadFile = File(...),
    module: ModuleKey = Form(...),
    title: Optional[str] = Form(default=None),
) -> UploadResponse:
    session = request.state.session
    upload_id = f"upload-{uuid4().hex[:8]}"
    filename = file.filename or "document"
    content = await file.read()
    stored_upload = upload_storage.store(
        upload_id=upload_id,
        filename=filename,
        content=content,
        content_type=file.content_type,
    )
    auth_backend.record_upload(
        upload_id=upload_id,
        workspace_id=session["workspace_id"],
        filename=filename,
        stored_path=stored_upload.stored_path,
        storage_backend=stored_upload.backend,
        storage_url=stored_upload.storage_url,
        content_type=file.content_type,
    )
    extracted_text = extract_text_from_upload(content, filename, file.content_type)
    source, artifact = store.ingest_upload(
        title=title or filename,
        module=module,
        source_type="upload",
        extracted_text=extracted_text,
    )
    return UploadResponse(
        upload=UploadRecord(
            id=upload_id,
            workspace_id=session["workspace_id"],
            filename=filename,
            stored_path=stored_upload.stored_path,
            storage_backend=stored_upload.backend,
            storage_url=stored_upload.storage_url,
            content_type=file.content_type,
            created_at=now_iso(),
        ),
        knowledge_source=source,
        artifact=artifact,
        message="Document uploaded and added to workspace memory.",
    )


@app.get("/api/uploads", response_model=list[UploadRecord])
def get_uploads(request: Request) -> list[UploadRecord]:
    session = request.state.session
    return [UploadRecord.model_validate(item) for item in auth_backend.list_uploads(session["workspace_id"])]


def request_last_touch_today() -> str:
    return now_iso().split("T", 1)[0]


def resolve_session(request: Request) -> Optional[dict[str, str]]:
    auth_mode = os.getenv("VXV_AUTH_MODE", "local").strip().lower()
    if auth_mode == "trusted-header":
        email_header = os.getenv("VXV_TRUSTED_EMAIL_HEADER", "X-Forwarded-Email")
        name_header = os.getenv("VXV_TRUSTED_NAME_HEADER", "X-Forwarded-User")
        role_header = os.getenv("VXV_TRUSTED_ROLE_HEADER", "X-Forwarded-Role")
        email = request.headers.get(email_header, "").strip()
        if not email:
            return None
        display_name = request.headers.get(name_header, "External User").strip() or "External User"
        role = request.headers.get(role_header, "member").strip() or "member"
        return auth_backend.create_trusted_session(
            email=email,
            display_name=display_name,
            workspace_id="workspace-vxv",
            role=role,
        )
    if auth_mode == "oidc":
        auth_header = request.headers.get("Authorization", "")
        token = auth_header.replace("Bearer ", "", 1).strip() if auth_header.startswith("Bearer ") else ""
        identity = oidc_verifier.validate(token)
        if identity is None:
            return None
        auth_backend.ensure_user(
            email=identity.email,
            workspace_id=identity.workspace_id,
            display_name=identity.display_name,
            role=identity.role,
        )
        return {
            "token": identity.token,
            "email": identity.email,
            "workspace_id": identity.workspace_id,
            "display_name": identity.display_name,
            "role": identity.role,
        }

    auth_header = request.headers.get("Authorization", "")
    token = auth_header.replace("Bearer ", "", 1).strip() if auth_header.startswith("Bearer ") else ""
    return auth_backend.get_session(token) if token else None


def extract_text_from_upload(content: bytes, filename: str, content_type: str | None) -> str:
    lowered_name = filename.lower()
    lowered_type = (content_type or "").lower()
    if lowered_name.endswith((".txt", ".md", ".json", ".csv")) or lowered_type.startswith("text/"):
        return trim_extracted_text(content.decode("utf-8", errors="ignore"))

    if lowered_name.endswith((".html", ".htm")) or lowered_type == "text/html":
        try:
            from bs4 import BeautifulSoup

            text = BeautifulSoup(content.decode("utf-8", errors="ignore"), "html.parser").get_text("\n").strip()
            return trim_extracted_text(text or "HTML uploaded successfully, but the extractor returned no content.")
        except Exception:
            return "HTML uploaded successfully. Text extraction is unavailable for this file in the current runtime."

    if lowered_name.endswith(".pdf") or lowered_type == "application/pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(BytesIO(content))
            text = "\n".join(page.extract_text() or "" for page in reader.pages).strip()
            return trim_extracted_text(text or "PDF uploaded successfully, but the text extractor returned no content.")
        except Exception:
            return "PDF uploaded successfully. Text extraction is unavailable for this file in the current runtime."

    if lowered_name.endswith(".docx") or lowered_type in {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    }:
        try:
            from docx import Document

            document = Document(BytesIO(content))
            text = "\n".join(paragraph.text for paragraph in document.paragraphs).strip()
            return trim_extracted_text(text or "DOCX uploaded successfully, but the extractor returned no content.")
        except Exception:
            return "DOCX uploaded successfully. Text extraction is unavailable for this file in the current runtime."

    if lowered_name.endswith(".pptx") or lowered_type in {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    }:
        try:
            from pptx import Presentation

            presentation = Presentation(BytesIO(content))
            slide_text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
            text = "\n".join(item for item in slide_text if item).strip()
            return trim_extracted_text(text or "PPTX uploaded successfully, but the extractor returned no content.")
        except Exception:
            return "PPTX uploaded successfully. Text extraction is unavailable for this file in the current runtime."

    if lowered_name.endswith(".xlsx") or lowered_type in {
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }:
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(filename=BytesIO(content), data_only=True)
            rows = []
            for sheet in workbook.worksheets:
                rows.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(value).strip() for value in row if value not in {None, ""}]
                    if values:
                        rows.append(" | ".join(values))
            text = "\n".join(rows).strip()
            return trim_extracted_text(text or "XLSX uploaded successfully, but the extractor returned no content.")
        except Exception:
            return "XLSX uploaded successfully. Text extraction is unavailable for this file in the current runtime."

    return (
        f"Uploaded file `{filename}` at {now_iso()}.\n"
        "The binary file is now attached to workspace memory, but it could not be converted into plain text automatically."
    )


def trim_extracted_text(value: str) -> str:
    limit = max(2000, int(os.getenv("VXV_INGEST_MAX_CHARS", "24000")))
    text = value.strip()
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + "\n\n[Truncated during ingestion to keep the workspace payload manageable.]"
